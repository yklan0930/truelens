# TrueLens 每日反馈汇总 · 第二步：读 _latest.json → 生成 Markdown 报告 + PPT 演示
# 用 python-pptx（已验证 1.0.2 在受管 venv）。缺失时自动 pip 安装；若仍不可用则仅出 md。
import json
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
REPORTS = ROOT / "docs" / "feedback-reports"
LATEST = REPORTS / "_latest.json"
STATE = REPORTS / ".last_run.txt"

# TrueLens 品牌色（同董事会报告）
NAVY = "0E2A47"
BLUE = "185FA5"
LIGHT = "EEF3FA"
GREY = "5B6B7B"


def ensure_pptx():
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.util import Inches, Pt, Emu  # noqa: F401
        from pptx.dml.color import RGBColor  # noqa: F401
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: F401

        return True
    except ImportError:
        try:
            subprocess.run(
                [sys.executable, "-m", "pip", "install", "--quiet", "python-pptx"],
                check=True,
            )
            return ensure_pptx()
        except Exception as e:  # noqa: BLE001
            print(f"WARN: python-pptx unavailable ({e}); skip PPT, keep Markdown.")
            return False


def fmt(ts):
    if not ts:
        return ""
    return str(ts)[:19].replace("T", " ")


def main():
    if not LATEST.exists():
        print("No _latest.json; run feedback_query.mjs first.")
        sys.exit(1)

    data = json.loads(LATEST.read_text(encoding="utf-8"))
    items = data.get("items", [])
    since = data.get("since", "")
    now = datetime.now(timezone.utc)
    local = now.astimezone()
    date_str = local.strftime("%Y-%m-%d")

    total = len(items)
    emoji = [i for i in items if i.get("type") == "emoji"]
    detailed = [i for i in items if i.get("type") == "detailed"]
    good = sum(1 for i in emoji if i.get("rating") == "good")
    bad = sum(1 for i in emoji if i.get("rating") == "bad")

    # ── Markdown 报告 ──
    md = []
    md.append(f"# TrueLens 用户反馈日报 · {date_str}\n")
    md.append(f"> 统计窗口：自 `{since}` 起的新反馈\n")
    if total == 0:
        md.append("\n## 本周期无新反馈 🎉\n")
        md.append("暂无用户提交反馈，继续观察。\n")
    else:
        md.append("\n## 关键指标\n")
        md.append(f"- 新增反馈总数：**{total}**")
        md.append(
            f"- 表情评分（👍 好 / 👎 差）：{good} / {bad}（共 {len(emoji)} 条）"
        )
        md.append(f"- 详细文字反馈：**{len(detailed)}** 条\n")

        if detailed:
            md.append("## 详细反馈\n")
            for idx, d in enumerate(detailed, 1):
                msg = (d.get("message") or "").strip() or "（无文字内容）"
                who = d.get("email") or d.get("userId") or "匿名"
                ctx = d.get("resultContext") or {}
                ctx_str = ""
                if isinstance(ctx, dict):
                    ai = ctx.get("aiProbability")
                    verdict = ctx.get("verdict")
                    fname = ctx.get("fileName")
                    if ai is not None:
                        ctx_str += f" ｜ AI味={ai}"
                    if verdict:
                        ctx_str += f" ｜ 判定={verdict}"
                    if fname:
                        ctx_str += f" ｜ 文件={fname}"
                md.append(
                    f"### {idx}. {who} · {fmt(d.get('createdAt'))}{ctx_str}\n\n> {msg}\n"
                )

        if emoji:
            md.append("## 表情评分明细\n")
            for e in emoji:
                who = e.get("email") or e.get("userId") or "匿名"
                r = (
                    "👍 好"
                    if e.get("rating") == "good"
                    else ("👎 差" if e.get("rating") == "bad" else str(e.get("rating")))
                )
                md.append(f"- {who} · {fmt(e.get('createdAt'))} · {r}")

    md_text = "\n".join(md)
    md_path = REPORTS / f"{date_str}.md"
    md_path.write_text(md_text, encoding="utf-8")

    # ── PPT 演示 ──
    pptx_ok = ensure_pptx()
    pptx_path = None
    if pptx_ok:
        pptx_path = build_pptx(
            date_str, total, good, bad, len(detailed), detailed, emoji, since
        )

    # ── 更新状态（下次窗口起点 = 本次生成时间）──
    STATE.write_text(now.isoformat(), encoding="utf-8")

    print(f"Report: {md_path}")
    if pptx_path:
        print(f"PPT:    {pptx_path}")
    print(f"Total={total} good={good} bad={bad} detailed={len(detailed)}")
    print(
        "SUMMARY:"
        + (" 本周期无新反馈" if total == 0 else f" {total} 条反馈（👍{good}/👎{bad}，详细 {len(detailed)}）")
    )


def build_pptx(date_str, total, good, bad, n_detailed, detailed, emoji, since):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    navy = RGBColor.from_string(NAVY)
    blue = RGBColor.from_string(BLUE)
    light = RGBColor.from_string(LIGHT)
    grey = RGBColor.from_string(GREY)
    white = RGBColor.from_string("FFFFFF")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def band(slide, title):
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = navy
        bar.line.fill.background()
        tf = bar.text_frame
        tf.margin_left = Inches(0.5)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = white

    # Slide 1 — 封面
    s = add_slide()
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = navy
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        ("TrueLens 用户反馈日报", 40, True, white),
        (date_str, 22, False, RGBColor.from_string("9DB8D6")),
        (
            "新增反馈 " + str(total) + " 条"
            if total
            else "本周期无新反馈 🎉",
            22,
            False,
            RGBColor.from_string("9DB8D6"),
        ),
    ]
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = col

    # Slide 2 — 关键指标
    s = add_slide()
    band(s, "关键指标")
    cards = [
        ("新增反馈", str(total), blue),
        ("👍 好评", str(good), RGBColor.from_string("1E8E5A")),
        ("👎 差评", str(bad), RGBColor.from_string("C0392B")),
        ("详细文字", str(n_detailed), grey),
    ]
    cw = Inches(2.9)
    gap = Inches(0.35)
    start_x = Inches(0.6)
    y = Inches(1.8)
    for i, (label, val, col) in enumerate(cards):
        x = start_x + i * (cw + gap)
        box = s.shapes.add_shape(1, x, y, cw, Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = light
        box.line.color.rgb = col
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(48)
        p1.font.bold = True
        p1.font.color.rgb = col
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(18)
        p2.font.color.rgb = grey
        p2.alignment = PP_ALIGN.CENTER
    note = s.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(12), Inches(1))
    note.text_frame.paragraphs[0].text = f"统计窗口：自 {since} 起的新反馈"
    note.text_frame.paragraphs[0].font.size = Pt(14)
    note.text_frame.paragraphs[0].font.color.rgb = grey

    # Slide 3 — 详细反馈
    s = add_slide()
    band(s, "详细反馈")
    tf = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8)).text_frame
    tf.word_wrap = True
    if detailed:
        for idx, d in enumerate(detailed[:8], 1):
            msg = (d.get("message") or "").strip() or "（无文字内容）"
            who = d.get("email") or d.get("userId") or "匿名"
            ctx = d.get("resultContext") or {}
            ctx_bits = []
            if isinstance(ctx, dict):
                if ctx.get("aiProbability") is not None:
                    ctx_bits.append(f"AI味={ctx.get('aiProbability')}")
                if ctx.get("verdict"):
                    ctx_bits.append(f"判定={ctx.get('verdict')}")
            ctx_str = (" ｜ " + "，".join(ctx_bits)) if ctx_bits else ""
            p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
            p.text = f"{idx}. {who} · {fmt(d.get('createdAt'))}{ctx_str}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = navy
            p2 = tf.add_paragraph()
            p2.text = "    " + msg
            p2.font.size = Pt(13)
            p2.font.color.rgb = grey
            p2.space_after = Pt(8)
    else:
        p = tf.paragraphs[0]
        p.text = "本周期无详细文字反馈。"
        p.font.size = Pt(16)
        p.font.color.rgb = grey

    # Slide 4 — 表情评分明细
    s = add_slide()
    band(s, "表情评分明细")
    tf = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8)).text_frame
    tf.word_wrap = True
    if emoji:
        for idx, e in enumerate(emoji[:12], 1):
            who = e.get("email") or e.get("userId") or "匿名"
            r = (
                "👍 好"
                if e.get("rating") == "good"
                else ("👎 差" if e.get("rating") == "bad" else str(e.get("rating")))
            )
            p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
            p.text = f"{idx}. {who} · {fmt(e.get('createdAt'))} · {r}"
            p.font.size = Pt(14)
            p.font.color.rgb = navy
            p.space_after = Pt(4)
    else:
        p = tf.paragraphs[0]
        p.text = "本周期无表情评分反馈。"
        p.font.size = Pt(16)
        p.font.color.rgb = grey

    out = REPORTS / f"{date_str}.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    main()
