# -*- coding: utf-8 -*-
"""
TrueLens 董事会汇报 deck builder.
Renders ONE shared design to both:
  - docs/TrueLens-Board-Report.pptx   (python-pptx)
  - docs/TrueLens-Board-Report.pdf    (reportlab canvas, 16:9 slide-like)
Run: python scripts/build_board_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.colors import HexColor
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

# ---------- palette ----------
NAVY  = '#0E2A47'
BLUE  = '#185FA5'
LBLUE = '#378ADD'
CYAN  = '#5BC0F8'
BG     = '#F3F7FC'
WHITE = '#FFFFFF'
DARK  = '#1B2838'
MUTED = '#5B6B7E'
GREEN = '#1E8E5A'
AMBER = '#C77B22'
RED   = '#C0392B'
CARD  = '#FFFFFF'
LINE  = '#D7E3F0'

W = 13.333
H = 7.5
PPTX_FONT = "Microsoft YaHei"
PDF_FONT  = "STSong-Light"

def _rgb(hexs):
    h = hexs.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# =====================================================================
#  Drawer abstraction (pptx + pdf share the same layout functions)
# =====================================================================
class PptxDrawer:
    def __init__(self, slide):
        self.s = slide
    def _box(self, x,y,w,h):
        return self.s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    def bg(self, color):
        self.rect(0,0,W,H,color)
    def rect(self, x,y,w,h, fill=None, line=None, line_w=1.0, radius=0.0):
        shp = self.s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius>0 else MSO_SHAPE.RECTANGLE,
            Inches(x),Inches(y),Inches(w),Inches(h))
        if fill: shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
        else: shp.fill.background()
        if line:
            shp.line.color.rgb = _rgb(line); shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        if radius>0:
            try:
                sp = shp.adjustments
                sp[0] = max(0.0, min(0.5, radius/ min(w,h)))
            except Exception: pass
        return shp
    def line(self, x1,y1,x2,y2, color, w=1.0):
        c = self.s.shapes.add_connector(2, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
        c.line.color.rgb = _rgb(color); c.line.width = Pt(w); c.shadow.inherit=False
        return c
    def text(self, x,y,w,h, text, size=14, color=DARK, bold=False,
             align='left', anchor='top', italic=False, spacing=1.15):
        tb = self._box(x,y,w,h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = {'top':MSO_ANCHOR.TOP,'middle':MSO_ANCHOR.MIDDLE,'bottom':MSO_ANCHOR.BOTTOM}[anchor]
        tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
        p = tf.paragraphs[0]; p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
        r = p.add_run(); r.text = text
        f = r.font; f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=PPTX_FONT
        f.color.rgb = _rgb(color)
        p.line_spacing = spacing
        return tb
    def bullets(self, x,y,w,h, items, size=14, color=DARK, gap=0.14,
                bullet='●', bullet_color=LBLUE, indent=0.30):
        for it in items:
            self.text(x, y, indent-0.06, size*1.4/72.0+0.05, bullet, size, bullet_color, bold=True)
            self.text(x+indent, y, w-indent, 0.6, it, size, color, spacing=1.12)
            y += gap + (size*1.12/72.0)

class PdfDrawer:
    def __init__(self, canvas):
        self.c = canvas
        self.W = W*72.0; self.H = H*72.0
        self.pt = 72.0
    def _y(self, y, h=0): return self.H - (y+h)*self.pt
    def bg(self, color):
        self.rect(0,0,W,H,color)
    def rect(self, x,y,w,h, fill=None, line=None, line_w=1.0, radius=0.0):
        xp=x*self.pt; yp=self._y(y,h); wp=w*self.pt; hp=h*self.pt
        if fill:
            self.c.setFillColor(HexColor(fill)); self.c.setStrokeColor(HexColor(fill))
        if line:
            self.c.setStrokeColor(HexColor(line)); self.c.setLineWidth(line_w)
        else:
            self.c.setStrokeColor(HexColor(fill or '#000000'))
        if radius>0:
            self.c.roundRect(xp,yp,wp,hp,radius*self.pt, stroke=1 if line else 0, fill=1 if fill else 0)
        else:
            self.c.rect(xp,yp,wp,hp, stroke=1 if line else 0, fill=1 if fill else 0)
        self.c.setLineWidth(1)
    def line(self, x1,y1,x2,y2, color, w=1.0):
        self.c.setStrokeColor(HexColor(color)); self.c.setLineWidth(w)
        self.c.line(x1*self.pt, self._y(y1), x2*self.pt, self._y(y2))
        self.c.setLineWidth(1)
    def _cw(self, ch, size):
        o=ord(ch)
        if ch==' ': return size*0.30
        if o<0x2E80 or (0xFF61<=o<=0xFF9F): return size*0.54
        return size*1.0
    def _wrap(self, text, maxw, size):
        lines=[]; cur=''; cw=0.0
        for ch in text:
            w=self._cw(ch,size)
            if cw+w>maxw and cur:
                lines.append(cur); cur=ch; cw=w
            else:
                cur+=ch; cw+=w
        if cur: lines.append(cur)
        return lines
    def text(self, x,y,w,h, text, size=14, color=DARK, bold=False,
             align='left', anchor='top', italic=False, spacing=1.18):
        self.c.setFillColor(HexColor(color))
        self.c.setFont(PDF_FONT, size)
        maxw = w*self.pt
        lines = self._wrap(text, maxw, size)
        lh = size*spacing
        total = lh*len(lines)
        if anchor=='top': top = self._y(y)
        elif anchor=='middle': top = self._y(y) - (total/2)
        else: top = self._y(y,h) + (h*self.pt - total)
        yy = top
        for ln in lines:
            if align=='left':
                self.c.drawString(x*self.pt, yy-lh+size*0.82, ln)
            elif align=='center':
                self.c.drawCentredString(x*self.pt + maxw/2, yy-lh+size*0.82, ln)
            else:
                self.c.drawRightString(x*self.pt + maxw, yy-lh+size*0.82, ln)
            yy -= lh
    def bullets(self, x,y,w,h, items, size=14, color=DARK, gap=0.14,
                bullet='●', bullet_color=LBLUE, indent=0.30):
        for it in items:
            self.c.setFillColor(HexColor(bullet_color)); self.c.setFont(PDF_FONT, size*0.8)
            self.c.drawString(x*self.pt, self._y(y)-size*0.25, bullet)
            lines = self._wrap(it, (w-indent)*self.pt, size)
            lh = size*1.12
            yy = self._y(y)
            for i,ln in enumerate(lines):
                self.c.setFillColor(HexColor(color)); self.c.setFont(PDF_FONT, size)
                self.c.drawString((x+indent)*self.pt, yy-size*0.25, ln)
                yy -= lh
            y += gap + (size*1.12/72.0)*max(1,len(lines))

# =====================================================================
#  shared chrome
# =====================================================================
PAGE = {'n':0}
def header(d, kicker, title):
    d.bg(BG)
    d.rect(0,0,W,1.2, NAVY)
    d.rect(0.5,0.42,0.13,0.42, LBLUE)
    d.text(0.78,0.36,9,0.3, kicker, 11.5, CYAN, bold=True)
    d.text(0.78,0.64,11.5,0.5, title, 25, WHITE, bold=True)
    # footer
    d.line(0.5,6.92,W-0.5,6.92, LINE, 0.75)
    d.text(0.5,7.0,8,0.34,"TrueLens · 董事会汇报 · 机密", 9, MUTED)
    PAGE['n']+=1
    d.text(W-1.4,7.0,0.9,0.34, f"{PAGE['n']:02d}", 9, MUTED, align='right')

def footerless(d):
    PAGE['n']=PAGE['n']  # cover/section don't count

# =====================================================================
#  slide layouts  (d = drawer, s = slide dict)
# =====================================================================
def s_cover(d, s):
    d.bg(NAVY)
    d.rect(0,0,0.35,H, LBLUE)
    d.rect(0.35,0,0.10,H, CYAN)
    # faint accent blocks
    d.rect(9.4,0.9,3.6,3.6, '#13335A', radius=0.4)
    d.rect(10.6,3.9,2.6,2.6, '#10304f', radius=0.4)
    d.text(1.1,1.5,11,0.5,"AI 内容真伪检测平台", 15, CYAN, bold=True)
    d.text(1.0,2.15,11.5,1.6,"TrueLens", 60, WHITE, bold=True)
    d.text(1.05,3.5,11,0.6,"发展回顾与战略展望 · 董事会汇报", 22, '#DCE9F7', bold=True)
    d.rect(1.05,4.35,2.2,0.06, LBLUE)
    d.text(1.05,4.6,11,0.5,"从「检测假图」到「内容真实性核验器」的演进之路", 14.5, '#A9C2DE')
    d.text(1.05,6.2,11,0.5,"2026 年 7 月   |   合伙人 Michael & 小毕   |   当前版本 v0.6.9", 12.5, '#8FB0D4')

def s_section(d, s):
    d.bg(NAVY)
    d.rect(0,0,W,H, NAVY)
    d.rect(0,3.0,0.35,H, LBLUE)
    d.text(1.1,2.55,2,1.2, s['num'], 70, '#1C4E86', bold=True)
    d.text(1.15,3.75,11,0.8, s['title'], 34, WHITE, bold=True)
    d.text(1.15,4.55,11,0.5, s['sub'], 15, CYAN)
    d.rect(1.15,5.35,2.0,0.05, LBLUE)

def s_agenda(d, s):
    header(d, "AGENDA", "汇报议程")
    items = s['items']
    colx = [0.7, 6.9]
    for i,it in enumerate(items):
        col = 0 if i<4 else 1
        row = i if i<4 else i-4
        x = colx[col]; y = 1.75 + row*1.18
        d.rect(x,y,5.6,1.0, CARD, line=LINE, line_w=1, radius=0.10)
        d.rect(x,y,0.12,1.0, LBLUE, radius=0.04)
        d.text(x+0.35,y+0.12,1.0,0.7, f"{i+1:02d}", 26, LBLUE, bold=True)
        d.text(x+1.25,y+0.16,4.2,0.45, it['t'], 15.5, DARK, bold=True)
        d.text(x+1.25,y+0.60,4.2,0.35, it['d'], 11.5, MUTED)

def s_exec(d, s):
    header(d, "EXECUTIVE SUMMARY", "执行摘要")
    d.text(0.7,1.55,12,0.5,"一句话定位", 13, BLUE, bold=True)
    d.rect(0.7,2.0,12.0,1.1, CARD, line=LINE, radius=0.10)
    d.rect(0.7,2.0,0.12,1.1, BLUE, radius=0.04)
    d.text(0.95,2.18,11.5,0.8,
        "TrueLens 是一个已上线（v0.6.8）的 AI 图片/视频真伪检测平台，"
        "正从「AI 假图检测器」升级为「内容真实性核验器」——以 Sightengine 检测 + C2PA 内容凭证验证双支柱，"
        "面向合规即服务与垂直摄影防伪两条可防守赛道。", 14.5, DARK, spacing=1.22)
    # KPI row
    kpis = s['kpis']
    kw=2.86; gap=0.22; x0=0.7; y=3.45
    for i,k in enumerate(kpis):
        x=x0+i*(kw+gap)
        d.rect(x,y,kw,2.55, CARD, line=LINE, radius=0.12)
        d.rect(x,y,kw,0.12, [BLUE,LBLUE,GREEN,AMBER][i%4], radius=0.04)
        d.text(x,y+0.45,kw,1.1,k['v'],40,[BLUE,LBLUE,GREEN,AMBER][i%4],bold=True,align='center')
        d.text(x,y+1.65,kw,0.4,k['k'],13.5,DARK,bold=True,align='center')
        d.text(x,y+2.05,kw,0.45,k['d'],10.5,MUTED,align='center')
    d.text(0.7,6.25,12,0.5,
        "关键判断：休闲 C 端「好玩」不会成为付费基础；真实且高增长的市场在 B2B 反欺诈与监管合规，"
        "而 2025–2026 全球监管集中落地正打开需求窗口。", 12.5, MUTED, italic=True)

def s_cards2(d, s):
    header(d, s['kicker'], s['title'])
    # left description
    d.text(0.7,1.6,5.6,0.5,s.get('lead',''),13.5,BLUE,bold=True)
    d.bullets(0.7,2.15,5.7,4.2, s['left'], size=13.5, gap=0.18)
    # right cards
    rx=6.7; rw=6.0
    cards=s['right']
    ch=(6.5-1.7)/len(cards) - 0.18
    y=1.75
    for c in cards:
        d.rect(rx,y,rw,ch, CARD, line=LINE, radius=0.10)
        d.rect(rx,y,0.12,ch, c.get('c',BLUE), radius=0.04)
        d.text(rx+0.32,y+0.12,rw-0.5,0.4,c['t'],14.5,DARK,bold=True)
        d.text(rx+0.32,y+0.55,rw-0.5,ch-0.6,c['d'],11.8,MUTED, spacing=1.12)
        y+=ch+0.18

def s_timeline(d, s):
    header(d, "HISTORY", "发展历程")
    ms = s['milestones']
    n=len(ms)
    x0=0.9; x1=12.4; y=3.55
    d.line(x0,y,x1,y, LINE, 2)
    step=(x1-x0)/(n-1)
    cols=[BLUE,LBLUE,GREEN,AMBER,RED,CYAN]
    for i,m in enumerate(ms):
        x=x0+i*step
        c=cols[i%len(cols)]
        d.rect(x-0.09,y-0.09,0.18,0.18,c,radius=0.09)
        above = (i%2==0)
        bx = x-1.05
        by = y-2.05 if above else y+0.30
        d.rect(bx,by,2.1,1.75, CARD, line=LINE, radius=0.10)
        d.rect(bx,by,2.1,0.42,c,radius=0.06)
        d.text(bx+0.12,by+0.06,1.9,0.32,m['v'],13,WHITE,bold=True)
        d.text(bx+0.12,by+0.52,1.9,1.15,m['t'],11.8,DARK,spacing=1.12)
        # connector
        d.line(x,y-(0.09 if above else -0.09), x if above else x, (by+1.75) if above else by, c, 1)
    d.text(0.7,6.35,12,0.5,"节奏：从「双模型误判」的试错，到 Sightengine 主引擎 100% 准确，再到 C2PA 验证把范式从「检测假」推向「证实真」。",12.5,MUTED,italic=True)

def s_tech(d, s):
    header(d, "TECHNOLOGY", "技术演进：检测引擎")
    stages=s['stages']
    x0=0.7; y=1.75; cw=2.95; gap=0.18
    cols=[RED,AMBER,BLUE,GREEN]
    for i,st in enumerate(stages):
        x=x0+i*(cw+gap)
        d.rect(x,y,cw,4.6, CARD, line=LINE, radius=0.12)
        d.rect(x,y,cw,0.95, cols[i%4], radius=0.08)
        d.text(x+0.18,y+0.12,cw-0.3,0.4,st['tag'],12,CYAN if i in(0,1) else WHITE,bold=True)
        d.text(x+0.18,y+0.48,cw-0.3,0.45,st['title'],15,WHITE,bold=True)
        d.bullets(x+0.20,y+1.15,cw-0.4,3.3, st['pts'], size=11.6, gap=0.16)
    d.rect(0.7,6.55,12.0,0.7, '#EAF2FB', line=LINE, radius=0.08)
    d.text(0.9,6.62,11.7,0.55,
        "视频路径：浏览器抽 8 帧 → 逐帧 Sightengine 图检测（串行 +1.1s 间隔 + 429 重试），免费层即可覆盖约 95% AI 生成视频。", 12, DARK)

def s_pricing(d, s):
    header(d, "BUSINESS MODEL", "产品与商业模式")
    plans=s['plans']
    x0=0.7; y=1.7; cw=2.92; gap=0.14
    hi=1  # Pro highlighted
    for i,p in enumerate(plans):
        x=x0+i*(cw+gap)
        hl = (i==hi)
        d.rect(x,y,cw,3.5, WHITE if hl else CARD, line=(LBLUE if hl else LINE), line_w=(2.0 if hl else 1.0), radius=0.12)
        if hl: d.rect(x,y,cw,0.5,LBLUE,radius=0.08); d.text(x,y+0.06,cw,0.4,"推荐",13,WHITE,bold=True,align='center')
        d.text(x,y+0.62,cw,0.5,p['name'],18,DARK if not hl else BLUE,bold=True,align='center')
        d.text(x,y+1.25,cw,0.8,p['price'],27,BLUE,bold=True,align='center')
        d.text(x,y+2.05,cw,0.35,p['credits'],12,MUTED,align='center')
        d.line(x+0.4,y+2.5,cw-0.8,y+2.5,LINE,0.75)
        d.bullets(x+0.30,y+2.62,cw-0.6,0.9, p['feat'], size=10.8, gap=0.10, bullet='—')
    # bottom strip
    d.rect(0.7,5.45,12.0,1.65, '#EAF2FB', line=LINE, radius=0.10)
    d.text(0.95,5.55,11.6,0.4,"支付与计费架构",13.5,BLUE,bold=True)
    d.bullets(0.95,6.0,11.7,1.0, s['notes'], size=12, gap=0.12, bullet='●', bullet_color=GREEN)

def s_status(d, s):
    header(d, "CURRENT STATUS", "当前进展")
    left=s['done']; right=s['todo']
    # done
    d.rect(0.7,1.7,5.9,4.9, CARD, line=LINE, radius=0.12)
    d.rect(0.7,1.7,5.9,0.6, GREEN, radius=0.08)
    d.text(0.95,1.78,5.5,0.42,"已达成 / Done",15,WHITE,bold=True)
    d.bullets(0.95,2.5,5.5,4.0, left, size=13, gap=0.22, bullet='✔', bullet_color=GREEN)
    # todo
    d.rect(6.7,1.7,5.9,4.9, CARD, line=LINE, radius=0.12)
    d.rect(6.7,1.7,5.9,0.6, AMBER, radius=0.08)
    d.text(6.95,1.78,5.5,0.42,"进行中 / 待办",15,WHITE,bold=True)
    d.bullets(6.95,2.5,5.5,4.0, right, size=13, gap=0.22, bullet='➜', bullet_color=AMBER)

def s_challenges(d, s):
    header(d, "CHALLENGES", "挑战与风险")
    items=s['items']
    x0=0.7; y=1.75; cw=3.85; gap=0.2; rh=2.25
    cols=[RED,AMBER,BLUE]
    for i,it in enumerate(items):
        col=i%3; row=i//3
        x=x0+col*(cw+gap); yy=y+row*(rh+0.2)
        c=cols[col%3]
        d.rect(x,yy,cw,rh, CARD, line=LINE, radius=0.12)
        d.rect(x,yy,cw,0.5,c,radius=0.07)
        d.text(x+0.18,yy+0.08,cw-0.3,0.36,it['t'],13.5,WHITE,bold=True)
        d.bullets(x+0.22,yy+0.66,cw-0.45,rh-0.8, it['pts'], size=11.6, gap=0.14)

def s_future(d, s):
    header(d, "FUTURE", "未来方向")
    items=s['items']
    x0=0.7; y=1.8; cw=3.85; gap=0.2; rh=2.2
    cols=[BLUE,LBLUE,GREEN,AMBER,RED,CYAN]
    for i,it in enumerate(items):
        col=i%3; row=i//3
        x=x0+col*(cw+gap); yy=y+row*(rh+0.18)
        c=cols[i%len(cols)]
        d.rect(x,yy,cw,rh, CARD, line=LINE, radius=0.12)
        d.rect(x,yy,0.12,rh,c,radius=0.04)
        d.text(x+0.3,yy+0.16,cw-0.5,0.45,f"{i+1}. {it['t']}",14.5,DARK,bold=True)
        d.text(x+0.3,yy+0.66,cw-0.55,rh-0.8,it['d'],11.8,MUTED,spacing=1.14)

def s_positioning(d, s):
    header(d, "STRATEGY", "战略落点：商业模式升级")
    # before -> after
    d.rect(0.7,1.75,5.7,2.1, '#FBEEEC', line='#F1C9C2', radius=0.12)
    d.text(0.95,1.85,5.3,0.4,"之前 · 单一定位",13,RED,bold=True)
    d.text(0.95,2.35,5.3,1.4,"「AI 假图检测器」\n\n靠单点检测赢，但生成模型持续进步，“检测假”是必输的仗。",13.5,DARK,spacing=1.2)
    d.rect(6.9,1.75,5.7,2.1, '#E7F3EC', line='#BFE2CC', radius=0.12)
    d.text(7.15,1.85,5.3,0.4,"升级 · 双支柱",13,GREEN,bold=True)
    d.text(7.15,2.35,5.3,1.4,"「内容真实性核验器」\n\n检测假（抓 AI）+ 验证真（C2PA 证真），在源头确权。",13.5,DARK,spacing=1.2)
    # three don'ts
    d.text(0.7,4.15,12,0.4,"三不碰（边界纪律）",13.5,BLUE,bold=True)
    d.bullets(0.7,4.6,12,1.0, s['donts'], size=13, gap=0.14, bullet='✕', bullet_color=RED)
    d.rect(0.7,5.75,12.0,1.35,'#EAF2FB',line=LINE,radius=0.10)
    d.text(0.95,5.85,11.6,0.4,"价值主张",13,BLUE,bold=True)
    d.text(0.95,6.28,11.6,0.7,s['value'],13,DARK,spacing=1.15)

def s_roadmap(d, s):
    header(d, "ACTION PLAN", "行动建议 / Roadmap")
    rows=s['rows']
    y=1.8
    pcol={'P0':RED,'P1':AMBER,'P2':BLUE,'P3':MUTED}
    for r in rows:
        c=pcol.get(r['p'],BLUE)
        d.rect(0.7,y,12.0,1.0, CARD, line=LINE, radius=0.10)
        d.rect(0.7,y,1.25,1.0, c, radius=0.06)
        d.text(0.7,y+0.30,1.25,0.5,r['p'],16,WHITE,bold=True,align='center')
        d.text(2.1,y+0.13,3.6,0.8,r['t'],14.5,DARK,bold=True,anchor='middle')
        d.text(5.9,y+0.13,6.6,0.8,r['d'],12,MUTED,anchor='middle',spacing=1.12)
        y+=1.12

def s_closing(d, s):
    d.bg(NAVY)
    d.rect(0,0,0.35,H, LBLUE)
    d.rect(1.1,2.2,2.0,0.06, LBLUE)
    d.text(1.1,2.5,11,1.2,"谢谢", 52, WHITE, bold=True)
    d.text(1.15,3.7,11,0.8,"以真实性为锚，做内容可信的基础设施。", 20, '#DCE9F7', bold=True)
    d.text(1.15,4.7,11,0.6,s['vision'], 14.5, '#A9C2DE', spacing=1.2)
    d.rect(1.15,5.7,11.0,0.05, '#1C4E86')
    d.text(1.15,5.95,11,0.5,"TrueLens  ·  truelens.top  ·  Michael & 小毕  ·  2026.07", 13, '#8FB0D4')

# =====================================================================
#  content
# =====================================================================
SLIDES = [
 {"type":"cover"},
 {"type":"section","num":"01","title":"执行摘要与议程","sub":"Executive Summary & Agenda"},
 {"type":"agenda","items":[
    {"t":"执行摘要","d":"定位、关键指标与核心判断"},
    {"t":"发展历程","d":"从双模型试错到 C2PA 验证"},
    {"t":"技术演进","d":"检测引擎与视频路径"},
    {"t":"产品与商业模式","d":"credits 计费与支付分层"},
    {"t":"进展与机遇","d":"上线状态、市场与监管"},
    {"t":"挑战·未来·行动","d":"风险、方向、路线图"},
 ]},
 {"type":"exec","kpis":[
    {"v":"v0.6.9","k":"产品版本","d":"已上线 v0.6.8 + C2PA 原型"},
    {"v":"100%","k":"17 图 QA 准确率","d":"AI / Real F1 双双 100%"},
    {"v":"$0.6–1.6B","k":"2025 全球市场","d":"CAGR 14–35%，2030s 冲 $10B+"},
    {"v":"~98%","k":"测算毛利率","d":"500 付费用户即打满成本"},
 ]},
 {"type":"section","num":"02","title":"发展历程","sub":"From trial-and-error to proven accuracy"},
 {"type":"timeline","milestones":[
    {"v":"v0.1–0.5","t":"立项与双模型探索：ViT 双模型加权，真实照片被严重误判"},
    {"v":"v0.6.0","t":"真实安全策略：仅显式水印判 AI；真实照零误判，但 AI 召回≈0%"},
    {"v":"v0.6.6","t":"配额月度化：credits 替代每日限额，奠定付费基础"},
    {"v":"v0.6.7","t":"修复「AI 味」百分比插值缺失 + PDF 分页优化"},
    {"v":"v0.6.8","t":"Layer A 支付基建上线：credits 余额 + 模拟收银台 + 套餐页"},
    {"v":"v0.6.9","t":"C2PA 内容凭证原型：从「检测假」迈向「验证真」"},
 ]},
 {"type":"section","num":"03","title":"技术与产品","sub":"Technology & Product"},
 {"type":"tech","stages":[
    {"tag":"阶段 1 · 失败","title":"双 HF 模型","pts":["Ateeqq + dima806 加权","Ateeqq 近二值，对人脸/真实照误判 ~100%","dima806 免费端点已死(400)","结论：双模型交叉验证失效"]},
    {"tag":"阶段 2 · 妥协","title":"真实安全策略","pts":["仅显式 AI 水印才判 likely_ai","其余诚实 likely_uncertain","Real F1=100%（零冤枉真人）","但 AI 召回≈0%（模型瓶颈）"]},
    {"tag":"阶段 3 · 突破","title":"Sightengine genai","pts":["图片/视频共用账户，作为主引擎","≥70→AI / ≤30→真实 / 其余弃权","17 图实测 100% 准、零误判","OCR 水印并行增强（中英文）"]},
    {"tag":"阶段 4 · 范式","title":"C2PA 验证","pts":["新增内容凭证验证端","双支柱：检测假 + 证实真","离线可验证、开源库可落地","已 build-passing（v0.6.9）"]},
 ]},
 {"type":"pricing","plans":[
    {"name":"Free","price":"¥0","credits":"3 次/月 高精度","feat":["基础引擎兜底","分享/PDF 导出","适合尝鲜"]},
    {"name":"Pro","price":"¥39","credits":"500 credits/月","feat":["图片 1 / 视频 8 credit","详细证据报告","优先级支持"]},
    {"name":"Business","price":"¥199","credits":"5000 credits/月","feat":["高用量团队","API 接入预留","定制支持"]},
    {"name":"加油包","price":"¥19","credits":"200 credits","feat":["随用随充","不绑定月度","成本可控"]},
 ],"notes":[
    "Layer A 已完成：credits 余额 + 原子扣减 + 全局月度 ops 预算闸 + 模拟收银台 + 后台充值 + /pricing 套餐页",
    "Layer B 待接：微信/支付宝真实收单（需企业商户号 + truelens.top ICP 备案）",
    "盈亏平衡约 15–20 付费用户；毛利测算约 98%（Starter $29≈¥210 含 1 万 ops）",
 ]},
 {"type":"section","num":"04","title":"进展 · 机遇 · 挑战","sub":"Status, Market & Challenges"},
 {"type":"status","done":[
    "已上线 v0.6.8（页脚版本号 + git SHA 可验证部署）",
    "17 图 QA：准确率 100%，AI / Real F1 双双 100%",
    "市场调研完成（规模 / 融资 / 监管 / 诈骗损失 / 付费意愿）",
    "临港 OPC 调研完成（15% 税惠 + AI 补贴 + 跨境金融）",
    "C2PA 原型代码 build-passing（v0.6.9，待 push）",
    "商业模式定位升级方案已论证（检测假 → 验证真）",
 ],"todo":[
    "git push v0.6.9 触发 Vercel 部署 C2PA 原型",
    "Layer B：接入微信/支付宝（商户号 + ICP 备案）",
    "摄影 GTM：比赛防伪落地页 + 白皮书",
    "临港 OPC 政策认定咨询（税惠 + 补贴申领）",
    "C2PA 验证端强化 → 合规即服务 API",
 ]},
 {"type":"cards2","kicker":"MARKET & REGULATION","title":"市场与监管机遇",
    "lead":"钱在 B2B 反欺诈与监管合规，不在个人「好玩」",
    "left":[
        "市场真实且高增长：2025 全球 $0.6–1.6B，CAGR 14–35%，2030 年代初冲 $10B+",
        "机构真金白银背书：Reality Defender 融资 $52.4M、GetReal $17.5M（Hany Farid / In-Q-Tel / Cisco）",
        "范式转移：生成模型持续进步，「检测假」必输；「证明真」(C2PA) 在源头确权是 2026 主潮",
        "休闲 C 端好奇不会成为付费基础——小 OPC 应避开通用企业反欺诈红海",
    ],
    "right":[
        {"t":"🇨🇳 中国标识办法","d":"《人工智能生成合成内容标识办法》+ 强制国标 GB45438-2025，已于 2025-09-01 施行（显式+隐式双标识、全链条责任）","c":RED},
        {"t":"🇪🇺 EU AI Act","d":"第 50 条 2026-08 生效，点名 C2PA 为满足披露的技术机制——出海企业合规刚需","c":BLUE},
        {"t":"🇺🇸 加州 SB 942","d":"2026-01-01 生效，要求合成内容提供披露工具","c":AMBER},
        {"t":"监管 = 加速器","d":"2025-09 → 2026-08 是 B2B 合规需求窗口，TrueLens 可吃 deadline 红利","c":GREEN},
    ]},
 {"type":"challenges","items":[
    {"t":"模型依赖","pts":["Sightengine 为外币(USD)支出，需持续跟随生成器迭代","自研对冲成本高，短期不替换主力"]},
    {"t":"支付 Layer B","pts":["需微信/支付宝企业商户号","需 truelens.top ICP 备案，规模化或需 EDI 许可证"]},
    {"t":"C2PA 工程坑","pts":["npm 安装门禁拦掉原生二进制下载","二进制来自 GitHub，中国网络可能拉不动"]},
    {"t":"资质门槛","pts":["司法鉴定/法院采信报告需许可+CNAS+≥3 鉴定人","小 OPC 不应争做——走「证明真」绕开"]},
    {"t":"竞争壁垒","pts":["Reality Defender / GetReal 资本与人才壁垒深","避开红海，切合规即服务 + 摄影 niche"]},
    {"t":"跨境数据","pts":["用户图片传 Sightengine（境外处理）需关注合规","临港有数据跨境绿色通道可用"]},
 ]},
 {"type":"section","num":"05","title":"未来与行动","sub":"Future & Action Plan"},
 {"type":"future","items":[
    {"t":"摄影 GTM","d":"比赛防伪 / 原创性证明；PRNU 相机指纹作差异化护城河，属「证明真」软性证明"},
    {"t":"临港政策申领","d":"申请 OPC「超级个体」认定，申领 15% 税惠 + Sightengine API 费用补贴"},
    {"t":"Layer B 收单","d":"接入微信/支付宝真实收单，把模拟收银台换成真钱闭环"},
    {"t":"C2PA 合规即服务","d":"中文 + AI 检测混合验证 API，帮平台/出海企业满足标识办法 + EU AI Act"},
    {"t":"自研检测器对冲","d":"频域(DCT/FFT)等方法作战略对冲，长期降低对 Sightengine 依赖"},
    {"t":"定位升级","d":"从「AI 假图检测器」升级为「内容真实性核验器」双支柱"},
 ]},
 {"type":"positioning","donts":[
        "不碰司法鉴定资质红线（许可 + CNAS + 鉴定人），用「证明真」绕开",
        "不自己当 CA / 签发端（证书基建重、合规重、被大厂垄断）",
        "不正面刚通用企业反欺诈红海（Reality Defender / GetReal 资本壁垒）",
    ],
    "value":"“我们帮你证明这份内容是否可信，并满足合规要求”——检测假 + 验证真，做内容真实性基础设施而非单一工具。"},
 {"type":"roadmap","rows":[
    {"p":"P0","t":"上线与政策","d":"git push v0.6.9；向临港管委会咨询 OPC 认定 + 15% 税惠 + Sightengine 补贴"},
    {"p":"P1","t":"真实收单","d":"拿微信/支付宝企业商户号 + ICP 备案，接入 Layer B 替换模拟收银台"},
    {"p":"P2","t":"摄影 GTM","d":"比赛防伪落地页 + 白皮书；用 PRNU 相机指纹建立护城河"},
    {"p":"P3","t":"合规即服务","d":"C2PA 验证端强化 → 合规 API；评估自研检测器作战略对冲"},
 ]},
 {"type":"closing","vision":"下一步：把 v0.6.9 推上线，启动临港政策申领，并按 P0→P3 顺序推进真实收单与垂直 GTM。"},
]

DISPATCH = {
    "cover":s_cover, "section":s_section, "agenda":s_agenda, "exec":s_exec,
    "cards2":s_cards2, "timeline":s_timeline, "tech":s_tech, "pricing":s_pricing,
    "status":s_status, "challenges":s_challenges, "future":s_future,
    "positioning":s_positioning, "roadmap":s_roadmap, "closing":s_closing,
}

# =====================================================================
#  render
# =====================================================================
def build_pptx(path):
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        d = PptxDrawer(slide)
        DISPATCH[s["type"]](d, s)
    prs.save(path)

def build_pdf(path):
    pdfmetrics.registerFont(UnicodeCIDFont(PDF_FONT))
    c = rl_canvas.Canvas(path, pagesize=(W*72.0, H*72.0))
    c.setTitle("TrueLens 董事会汇报")
    for s in SLIDES:
        d = PdfDrawer(c)
        DISPATCH[s["type"]](d, s)
        c.showPage()
    c.save()

if __name__ == "__main__":
    out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
    os.makedirs(out, exist_ok=True)
    pptx = os.path.join(out, "TrueLens-Board-Report.pptx")
    pdf  = os.path.join(out, "TrueLens-Board-Report.pdf")
    build_pptx(pptx)
    build_pdf(pdf)
    print("PPTX ->", pptx)
    print("PDF  ->", pdf)
